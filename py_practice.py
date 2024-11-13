from __future__ import absolute_import
from __future__ import print_function
#from .project import ProjectBase
#import svtools.common.baseaccess as baseaccess
#import namednodes
#from namednodes.precondition import MultipleAccesses
import os
import time
import re
from math import log2
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
import pandas as pd
import dataframe_image as dfi
# debug branch
        
def test1():
    
    path = r'C:\Scripts'
    f = open(F"{path}\\B0_PO_D4388HN400009_Q4V3_MRDIMM_Hynix_Rea_2Rx4_8800_CH0_Verbose.log", "r")
    csvf = open(F"{path}\\f.xlsx", "w")
    csvf.write(F"Channel,DIMM,RCS\n")
    lines = f.readlines()
    idx = 0
    text = "behind"
    start_pattern = "START_SOCKET_(\d)_DIMMINFO_TABLE"
    stop_pattern = "STOP_SOCKET_(\d)_DIMMINFO_TABLE"
    analysis = 0
    ch = []
    dimm = []
    rcd = []
    for line in lines:
        if text in line:
            #index = line.index(text)
            #print (F"{text} at index {index} in line {idx} of", line)
            words = line.split()
            words_len = len(words)
            match = 0
            for count in range(words_len):
                #match = re.search(r"DDR-(\d{4})", words[count])
                if re.search("DDR-(\d{4})", words[count]):
                    word_split = words[count].split("-")
                    #print (F"Speed: {word_split[1]}")
        elif re.match(start_pattern, line):
            analysis = 1
        elif re.match(stop_pattern, line):
            analysis = 0
            for count in range(0, 12):
                csvf.write(F"{ch[count]},{dimm[count]},{rcd[count]}\n")
            csvf.close()
            continue
        elif analysis == 1:
            if re.match("S", line):
                words = line.split("|")
                words_len = len(words) - 2
                for count in range(1, words_len+1):
                    words[count] = words[count].replace("Channel ", "")
                    ch.append(words[count])
            if re.match("0", line):
                words = line.split("|")
                words_len = len(words) - 2
                for count in range(1, words_len+1):
                    #print (words[count])
                    dimm_vendor = words[count]
                    if re.search("DIMM", words[count]):
                        dimm_vendor = words[count].replace('DIMM:', '')
                    dimm_vendor = dimm_vendor.lstrip()
                    dimm_vendor = dimm_vendor.rstrip()
                    dimm.append(dimm_vendor)
            if re.match(" ", line) and re.search("RCD", line) and (not re.search("Rev", line)):
                words = line.split("|")
                words_len = len(words) - 1
                for count in range(1, words_len):
                    #print (words[count])
                    rcd_vendor = words[count]
                    if re.search("RCD", words[count]):
                        rcd_vendor = words[count].replace('RCD:', '')
                    rcd_vendor = rcd_vendor.lstrip()
                    rcd_vendor = rcd_vendor.rstrip()
                    rcd.append(rcd_vendor)
        idx = idx + 1
    csvf.close()
    
    #slide_layouts[8] can insert picture
    ppt = Presentation()
    title_slide = ppt.slide_layouts[6]
    section_slide = ppt.slide_layouts[1]
    blank_slide = ppt.slide_layouts[5]
    title_image_slide = ppt.slide_layouts[4]
    end_slide = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(title_slide)
    #slide.shapes.placeholders[0].text = 'UPM Calculation'
    #slide.shapes.placeholders[1].text = 'Brandon Guo'
    #picture = slide.placeholders[13].insert_picture(os.path.join(img_path, pic))
    #slide = ppt.slides.add_slide(title_image_slide)
    index = 0
    #for pHolder in slide.shapes: #placeholders:
    #    if pHolder.is_placeholder:
    #        print (pHolder.placeholder_format.type)
    #picture = slide.placeholders[0].insert_picture('sun.png')
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(2)
    height = Inches(2)
    picture = slide.shapes.add_picture("sun.png",left,top,width,height)
       
    txBox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Brandon Guo"
    p.font.size = Pt(40)
    p.alignment = 2
        
    table_placeholder = slide.shapes.add_table(3, 4, Inches(5), Inches(4), Inches(5), Inches(3))
    table = table_placeholder.table
    rows=3
    cols=4
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = f"Cell {row_idx + 1}, {col_idx + 1}"
    
    data_file = F"{path}\\ff.xlsx"
    wb = load_workbook(data_file)
    print("Found the following worksheets:")
    for sheetname in wb.sheetnames:
        print(sheetname)
    ws = wb['f']
    all_rows = list(ws.rows)
    row_len = len(all_rows)
    print(row_len)
    for idx in range(row_len):
        count = 0
        for cell in all_rows[idx]:
            cell_len = len(all_rows[idx])
            if count != cell_len - 1:
                print(cell.value, end = " ")
            else:
                print(cell.value)
            count = count + 1
            
    df = pd.read_excel(data_file)
    #pivot1 = pd.pivot_table(df,index='DIMM',values=['Channel'], columns='RCS', aggfunc=['sum','mean'], margins=True, margins_name='Total')
    pivot1 = pd.pivot_table(df,index='DIMM',values=['Channel'], columns='RCS', aggfunc=['sum','mean'])
    #dfi.export(table.style.applymap(_highlight_upm, subset=['UPM']), os.path.join(path, 'Slides', f"[Tb]_UPM.png"))
    print(pivot1)
    pivot1.style.map(_highlight_upm, subset=['sum','mean']).to_excel('DIMM.xlsx')
    dfi.export(pivot1.style.map(_highlight_upm, subset=['sum','mean']), f"DIMM.png")
    picture = slide.shapes.add_picture("DIMM.png",left,Inches(5),width,height)
    ppt.save(F"{path}\\f.ppt")
    
def _highlight_upm(s):
    if s < 6:
        return 'background-color:lightgreen'
    elif 5 < s < 10:
        return 'background-color:yellow'
    elif s > 10:
        return 'background-color:orangered'
    else:
        return None
    
    
    


        


